"""居宅介護支援事業所向け 介護保険資格確認等WEBサービス紹介チラシ (A4縦・両面・PowerPoint)
   v3: ケアマネ目線リニューアル — シーン中心 / Before・After で「便利になる」を伝える"""
import os
from pptx import Presentation
from pptx.util import Mm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# カラーパレット (Teal Trust + 暖色アクセント)
NAVY = RGBColor(0x06, 0x36, 0x5A)
TEAL = RGBColor(0x02, 0x80, 0x90)
SEAFOAM = RGBColor(0x00, 0xA8, 0x96)
MINT = RGBColor(0x88, 0xDB, 0xC2)
CREAM = RGBColor(0xFE, 0xF6, 0xEE)
WARM = RGBColor(0xFF, 0xF3, 0xD6)
ACCENT = RGBColor(0xE8, 0x6A, 0x0C)
RED_SOFT = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x10, 0x80, 0x52)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x29, 0x33)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_BG = RGBColor(0xF7, 0xF5, 0xF0)
LIGHT_BORDER = RGBColor(0xD1, 0xD5, 0xDB)
PINK_BG = RGBColor(0xFD, 0xEC, 0xEC)
GREEN_BG = RGBColor(0xE7, 0xF7, 0xEE)

JP_SANS = 'HG丸ゴシックM-PRO'
JP_POP = 'HGP創英角ﾎﾟｯﾌﾟ体'
JP_MIN = 'HG丸ゴシックM-PRO'

prs = Presentation()
prs.slide_width = Mm(210)
prs.slide_height = Mm(297)
blank = prs.slide_layouts[6]


def text(slide, x, y, w, h, content, *, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT,
         font=JP_SANS, anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    tb = slide.shapes.add_textbox(Mm(x), Mm(y), Mm(w), Mm(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Mm(0); tf.margin_right = Mm(0)
    tf.margin_top = Mm(0); tf.margin_bottom = Mm(0)
    tf.vertical_anchor = anchor
    lines = content.split('\n') if isinstance(content, str) else content
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def rect(slide, x, y, w, h, *, fill=WHITE, line=None, corner=None, line_w=0.75):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if corner is not None else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, Mm(x), Mm(y), Mm(w), Mm(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    if corner is not None and hasattr(s, 'adjustments'):
        s.adjustments[0] = corner
    if s.has_text_frame:
        s.text_frame.text = ''
    return s


def circle(slide, x, y, d, *, fill=NAVY, label='', size=18, color=WHITE, bold=True, font=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Mm(x), Mm(y), Mm(d), Mm(d))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    tf = s.text_frame
    tf.margin_left = Mm(0); tf.margin_right = Mm(0)
    tf.margin_top = Mm(0); tf.margin_bottom = Mm(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = font if font else JP_POP
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return s


def arrow(slide, x, y, w, h, *, fill=ACCENT):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Mm(x), Mm(y), Mm(w), Mm(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


# =====================================================================
# Slide 1 (表面) - シーン Before/After: ケアマネの「あるある」5 場面
# =====================================================================
s1 = prs.slides.add_slide(blank)

# ── ヘッダー ──
rect(s1, 0, 0, 210, 26, fill=NAVY)
text(s1, 15, 5, 180, 8, '介護保険資格確認等WEBサービス',
     size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=JP_POP)
text(s1, 15, 16, 180, 6, '令和8年(2026年)4月 開始  ・  居宅介護支援事業所さま向けご案内',
     size=10, color=MINT, align=PP_ALIGN.CENTER)

# ── ヒーロー ──
rect(s1, 0, 28, 210, 30, fill=CREAM)
rect(s1, 0, 28, 210, 1.5, fill=ACCENT)
rect(s1, 0, 56.5, 210, 1.5, fill=ACCENT)
text(s1, 15, 32, 180, 9, 'こんなこと、ありませんか?',
     size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER, font=JP_POP)
text(s1, 15, 44, 180, 7, '居宅ケアマネ業務 5 つの「あるある」場面が、新しい WEB サービスでこう変わります。',
     size=11, color=DARK, align=PP_ALIGN.CENTER)
text(s1, 15, 51, 180, 5, '〜 紙・FAX・電話・コピー回収 から、API による電子連携へ 〜',
     size=9, color=GRAY, align=PP_ALIGN.CENTER)

# ── 5 シーンカード ──
# 各カード: 180mm × 32mm
# レイアウト: [場面 badge 22mm] [タイトル+ビフォー+アフター 158mm]

scenes = [
    ('01', '月初の 利用票・提供票 配布',
     '「毎月10ヶ所のサービス事業所へ FAX 流すだけで半日…」',
     '紙でプリント → FAX 送信 → 受領確認の電話',
     '介護ソフトから 一斉電子送信 → 受信も自動確認',
     ACCENT),
    ('02', '認定更新時の 証コピー受領',
     '「認定通知書をコピーで受け取って、介護ソフトに手入力…」',
     'コピー受領 → 手入力 → ミスで返戻リスク',
     '被保険者番号で API 照会 → 最新情報を自動取得',
     TEAL),
    ('03', '毎年8月の 負担割合証 回収',
     '「8月になると35名分の負担割合証コピー回収に走り回る…」',
     '各家庭訪問 → コピー受領 → 介護ソフト更新',
     '訪問・受領 不要 → API で 自動反映',
     SEAFOAM),
    ('04', '新規受入れ時の 引継ぎ情報',
     '「前の居宅から十分な情報がもらえず、利用者にも聞き直し…」',
     '紙の引継ぎ書頼り / 情報不足 / 再ヒアリング',
     '介護情報基盤から 過去ケアプラン履歴 を参照',
     NAVY),
    ('05', '認定有効期限 の チェック',
     '「気付いたら期限切れで請求エラー…月末に大慌て…」',
     '介護ソフトを目視チェック / 個別カレンダー管理',
     '期限切れ前に 自動アラート → 計画的更新が可能',
     RED_SOFT),
]

scene_y_start = 62
card_h = 31
card_gap = 2

for i, (no, title, voice, before, after, color) in enumerate(scenes):
    cy = scene_y_start + i * (card_h + card_gap)

    # カード全体枠
    rect(s1, 15, cy, 180, card_h, fill=WHITE, line=LIGHT_BORDER, corner=0.05, line_w=0.5)

    # 左 番号 badge (22mm × card_h)
    rect(s1, 15, cy, 22, card_h, fill=color, corner=0.05)
    rect(s1, 22, cy, 15, card_h, fill=color)  # 角丸を打ち消すための補助
    text(s1, 15, cy + 4, 22, 10, no,
         size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=JP_POP)
    text(s1, 15, cy + 18, 22, 6, '場面',
         size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=JP_POP)

    # タイトル
    text(s1, 40, cy + 2.5, 152, 7, title,
         size=13, bold=True, color=NAVY, font=JP_POP)

    # 利用者ボイス (吹き出し風)
    text(s1, 40, cy + 9.5, 152, 5, voice,
         size=9, color=GRAY, font=JP_SANS)

    # Before/After 2 列横並び
    # Before (左, ピンク背景)
    rect(s1, 40, cy + 15, 73, 13, fill=PINK_BG, corner=0.1)
    text(s1, 42, cy + 16, 18, 5, '今まで',
         size=8, bold=True, color=RED_SOFT, font=JP_POP)
    text(s1, 42, cy + 21, 70, 6, before,
         size=9, color=DARK, line_spacing=1.15)

    # 矢印
    arrow(s1, 114, cy + 18, 6, 7, fill=ACCENT)

    # After (右, 緑背景)
    rect(s1, 122, cy + 15, 70, 13, fill=GREEN_BG, corner=0.1)
    text(s1, 124, cy + 16, 18, 5, 'これから',
         size=8, bold=True, color=GREEN, font=JP_POP)
    text(s1, 124, cy + 21, 67, 6, after,
         size=9, bold=True, color=NAVY, line_spacing=1.15)

# ── フッター誘導 ──
foot_y = 228
rect(s1, 0, foot_y, 210, 50, fill=NAVY)
text(s1, 15, foot_y + 5, 180, 7, 'つまり、ケアマネ業務はこう変わります',
     size=11, color=MINT, align=PP_ALIGN.CENTER)
text(s1, 15, foot_y + 14, 180, 11, '紙とFAXに追われる毎月から、',
     size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=JP_POP)
text(s1, 15, foot_y + 25, 180, 11, '利用者と向き合う時間 を取り戻せます。',
     size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, font=JP_POP)
text(s1, 15, foot_y + 40, 180, 6, '▶  裏面で 1ヶ月の業務 が どう変わるか をご紹介します',
     size=10, color=MINT, align=PP_ALIGN.CENTER)

# 注釈
text(s1, 15, 282, 180, 5,
     '※ 同意の有効期間・再取得要否は現時点での情報。令和8年4月の正式版で確定予定',
     size=7, color=GRAY, align=PP_ALIGN.CENTER)


# =====================================================================
# Slide 2 (裏面) - 1ヶ月のケアマネ業務 Before/After タイムライン + 行動
# =====================================================================
s2 = prs.slides.add_slide(blank)

# ── ヘッダー ──
rect(s2, 0, 0, 210, 26, fill=NAVY)
text(s2, 15, 5, 180, 8, 'ケアマネの 1ヶ月、こう変わる',
     size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=JP_POP)
text(s2, 15, 16, 180, 6, '〜 月初・月中・月末 業務の Before / After ・ 利用シーン別の便利ポイント 〜',
     size=10, color=MINT, align=PP_ALIGN.CENTER)

# ── 1ヶ月タイムライン ──
text(s2, 15, 30, 180, 7, '1ヶ月のケアマネ業務 Before / After',
     size=14, bold=True, color=NAVY, font=JP_POP, align=PP_ALIGN.CENTER)
text(s2, 15, 37, 180, 5, '— 現状の手間と、新サービス導入後の姿を時系列で比較 —',
     size=9, color=GRAY, align=PP_ALIGN.CENTER)

timeline = [
    ('月  初', '1-3日',
     '利用票・提供票 を サービス事業所へFAX  /  受領確認の電話',
     '介護ソフトから 電子送信  /  受信確認も自動  → 半日の作業が 30 分に',
     ACCENT),
    ('月  中', '10-20日',
     '認定更新・期限管理を 目視チェック  /  証コピー受領',
     '認定情報を API 自動取得  /  期限切れ前に 自動アラート',
     TEAL),
    ('月  末', '27-31日',
     '翌月利用票 作成のため 各事業所と電話・FAX で予定すり合わせ',
     '前月実績データを そのまま活用  /  電子で計画調整',
     SEAFOAM),
    ('年1回', '8月',
     '負担割合証コピー を 全利用者から回収  (担当35件 × 訪問)',
     '訪問・回収 不要  /  API で 全員分 一括自動反映',
     NAVY),
]

tl_y = 44
for i, (period, days, before, after, color) in enumerate(timeline):
    row_y = tl_y + i * 26
    # 期間ラベル
    rect(s2, 15, row_y, 28, 22, fill=color, corner=0.08)
    text(s2, 15, row_y + 3, 28, 8, period,
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=JP_POP)
    text(s2, 15, row_y + 13, 28, 6, days,
         size=9, color=WHITE, align=PP_ALIGN.CENTER)

    # Before
    rect(s2, 46, row_y, 71, 22, fill=PINK_BG, corner=0.08)
    text(s2, 49, row_y + 1.5, 14, 5, '今まで',
         size=8, bold=True, color=RED_SOFT, font=JP_POP)
    text(s2, 49, row_y + 7, 67, 14, before,
         size=9, color=DARK, line_spacing=1.3)

    # Arrow
    arrow(s2, 119, row_y + 7, 6, 8, fill=ACCENT)

    # After
    rect(s2, 127, row_y, 68, 22, fill=GREEN_BG, corner=0.08)
    text(s2, 130, row_y + 1.5, 14, 5, 'これから',
         size=8, bold=True, color=GREEN, font=JP_POP)
    text(s2, 130, row_y + 7, 64, 14, after,
         size=9, bold=True, color=NAVY, line_spacing=1.3)

# ── 削減効果サマリ (数値で訴求) ──
sum_y = 152
rect(s2, 15, sum_y, 180, 28, fill=WARM, corner=0.05, line=ACCENT, line_w=1.2)
text(s2, 20, sum_y + 3, 170, 7, '【期待される効果】',
     size=11, bold=True, color=ACCENT, font=JP_POP)

# 3 指標を横並び
metrics = [
    ('FAX 業務', '大幅減', '計画書・利用票・提供票の電子化'),
    ('コピー受領', 'ほぼゼロ', '被保険者証 / 負担割合証 / 認定通知'),
    ('期限切れ', '0 件 へ', '介護ソフト 自動アラートで未然防止'),
]
for i, (label, val, sub) in enumerate(metrics):
    mx = 20 + i * 58
    text(s2, mx, sum_y + 10, 56, 5, label,
         size=10, color=DARK, align=PP_ALIGN.CENTER)
    text(s2, mx, sum_y + 15, 56, 7, val,
         size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, font=JP_POP)
    text(s2, mx, sum_y + 23, 56, 4, sub,
         size=7, color=GRAY, align=PP_ALIGN.CENTER)

# ── 利用者にとっての安心ポイント ──
ben_y = 184
text(s2, 15, ben_y, 180, 7, 'おまけ: 利用者さま側のメリットも',
     size=12, bold=True, color=NAVY, font=JP_POP)
text(s2, 15, ben_y + 6, 180, 5, '— ケアマネからの「コピーちょうだい」「印鑑お願いします」が減ります —',
     size=9, color=GRAY)

user_benefits = [
    ('🛡', '同意の範囲を 利用者自身が管理'),
    ('📎', '被保険者証 / 負担割合証コピーの 手渡し不要'),
    ('🔄', '事業所変更時も 過去の計画 を引き継ぎ可能'),
]
ben_x_start = 15
for i, (icon, t) in enumerate(user_benefits):
    bx = ben_x_start + i * 60
    rect(s2, bx, ben_y + 14, 56, 11, fill=LIGHT_BG, corner=0.1)
    text(s2, bx + 2, ben_y + 16, 8, 7, icon,
         size=14, align=PP_ALIGN.CENTER, font=JP_SANS)
    text(s2, bx + 11, ben_y + 18, 44, 6, t,
         size=9, color=DARK, line_spacing=1.2)

# ── ご利用までの準備 ──
prep_y = 211
rect(s2, 15, prep_y, 180, 32, fill=LIGHT_BG, corner=0.03)
text(s2, 22, prep_y + 3, 166, 6, 'ご利用までに必要な準備',
     size=12, bold=True, color=NAVY, font=JP_POP)

preps = [
    ('1', 'クライアント証明書', '介護DX証明書 / 介護保険証明書'),
    ('2', 'ユーザ ID・PW 払い出し', '事業所/管理者/一般 の3階層'),
    ('3', 'マイナ資格確認アプリ', 'PCカードリーダー or スマホNFC'),
    ('4', '対応介護ソフト', '各ベンダーが API 実装中'),
]
for i, (n, t, d) in enumerate(preps):
    pcx = 22 + i * 44
    circle(s2, pcx, prep_y + 11, 7, fill=NAVY, label=n, size=10, color=WHITE)
    text(s2, pcx + 9, prep_y + 11, 36, 5, t,
         size=9, bold=True, color=NAVY)
    text(s2, pcx + 9, prep_y + 17, 36, 6, d,
         size=7, color=GRAY, line_spacing=1.15)

# ── KT グループの提案 ──
prop_y = 247
rect(s2, 15, prop_y, 180, 22, fill=NAVY, corner=0.05)
text(s2, 22, prop_y + 3, 166, 5, 'KTグループからのご提案',
     size=9, color=MINT)
text(s2, 22, prop_y + 8, 166, 7, '居宅 ⇄ サービス事業所 の電子連携、まずは KTグループと。',
     size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=JP_POP)
text(s2, 22, prop_y + 16, 166, 5,
     '訪問介護・福祉用具事業所 と スムーズに連携できる受信体制を構築中。',
     size=9, color=MINT, align=PP_ALIGN.CENTER)

# ── 連絡先 ──
contact_y = 272
rect(s2, 15, contact_y, 180, 14, fill=WHITE, line=ACCENT, corner=0.08, line_w=1.2)
text(s2, 20, contact_y + 1.5, 170, 5, 'お問い合わせ / 連携体制のご相談',
     size=9, bold=True, color=ACCENT, font=JP_POP)
text(s2, 20, contact_y + 7, 170, 6,
     '株式会社ケイ・ティ・サービス / KTグループ      TEL: 043-XXX-XXXX      Web: kt-group.co.jp',
     size=10, color=DARK)

# ── 注釈 ──
text(s2, 15, 289, 180, 5,
     '※ 本資料は現時点の情報 (国保中央会公表の暫定版仕様) に基づき作成。同意有効期間等の運用詳細は令和8年4月の正式版で確定予定。',
     size=7, color=GRAY, align=PP_ALIGN.CENTER)


OUT = os.path.join(os.path.dirname(__file__), 'flyer_kaigo_web_service.pptx')
try:
    prs.save(OUT)
except PermissionError:
    OUT = os.path.join(os.path.dirname(__file__), 'flyer_kaigo_web_service_v3.pptx')
    prs.save(OUT)
print('Saved:', OUT)
